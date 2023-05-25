import collections
import collections.abc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from googletrans import Translator

print("Input the original and translate in .txt")
inpSinger = input("Singer: ")
inpTitle = input("Title: ")

prs = Presentation()

bullet_slide_layout = prs.slide_layouts[0]

with open('input_orig.txt', 'r', encoding="utf-8") as f:
    input_orig = f.read()

text_orig_list = input_orig.split("\n")
text_orig_list = [x for x in text_orig_list if x]  # delete empty elements

with open('input_trans.txt', 'r', encoding="utf-8") as f:
    input_trans = f.read()
text_trans_list = input_trans.split("\n")
text_trans_list = [x for x in text_trans_list if x]  # delete empty elements

if len(text_trans_list) != len(text_orig_list):
    print("Translate not exist, use auto translate")
    translator = Translator()
    text_trans_list = []
    for i in range(len(text_orig_list)):
        text = text_orig_list[i]
        translation = translator.translate(text, dest="ru")
        text_trans_list.append(translation.text)

for i in range(len(text_orig_list)):
    slide = prs.slides.add_slide(bullet_slide_layout)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    text_frame = slide.shapes[0].text_frame

    text_frame.text = text_orig_list[i]  # текст оригинала
    text_frame.paragraphs[0].font.name = 'TimesNewRoman'

    line = len(text_orig_list[i])
    if line <= 10:
        font_size = 90
    else:
        font_size = round(900 / line)

    text_frame.paragraphs[0].font.size = Pt(font_size)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    tf = body_shape.text_frame
    tf.text = text_trans_list[i]  # текст перевода

    line = len(text_trans_list[i])
    if line <= 10:
        font_size = 90
    else:
        font_size = round(900 / line)

    tf.paragraphs[0].font.name = 'TimesNewRoman'
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

prs.save(r"C:\Users\halat\Desktop\songs_pptx\{} - {}.pptx".format(inpSinger, inpTitle))

input("Success")